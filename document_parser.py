import os
from pathlib import Path
import pdfplumber
from pptx import Presentation
from docx import Document as DocxDocument

SUPPORTED_TEXT_TYPES = {'pdf', 'pptx', 'docx', 'txt'}


def extract_document_pages(file_path: str, file_type: str) -> list[dict]:
    file_type = file_type.lower()
    if file_type == 'pdf':
        return _extract_pdf_pages(file_path)
    if file_type == 'pptx':
        return _extract_pptx_slides(file_path)
    if file_type == 'docx':
        return _extract_docx_sections(file_path)
    if file_type == 'txt':
        return _extract_txt(file_path)
    raise ValueError(f'Unsupported file type: {file_type}')


def _extract_pdf_pages(file_path: str) -> list[dict]:
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for page_index, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ''
            pages.append({'page_number': page_index, 'text': text.strip(), 'source_type': 'pdf'})
    return pages


def _extract_pptx_slides(file_path: str) -> list[dict]:
    presentation = Presentation(file_path)
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                parts.append(shape.text)
        slides.append({'page_number': index, 'text': '\n'.join(parts).strip(), 'source_type': 'pptx'})
    return slides


def _extract_docx_sections(file_path: str) -> list[dict]:
    document = DocxDocument(file_path)
    sections = []
    current_text = []
    section_index = 1
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if len(current_text) >= 1000:
            sections.append({'page_number': section_index, 'text': ' '.join(current_text).strip(), 'source_type': 'docx'})
            current_text = []
            section_index += 1
        current_text.append(text)
    if current_text:
        sections.append({'page_number': section_index, 'text': ' '.join(current_text).strip(), 'source_type': 'docx'})
    if not sections:
        sections.append({'page_number': 1, 'text': '', 'source_type': 'docx'})
    return sections


def _extract_txt(file_path: str) -> list[dict]:
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as input_file:
        text = input_file.read().strip()
    return [{'page_number': 1, 'text': text, 'source_type': 'txt'}]


def chunk_text(text: str, max_tokens: int = 500, overlap: int = 100) -> list[str]:
    words = text.split()
    if not words:
        return []
    chunks = []
    start = 0
    while start < len(words):
        end = min(len(words), start + max_tokens)
        chunks.append(' '.join(words[start:end]))
        if end == len(words):
            break
        start = end - overlap
    return chunks
